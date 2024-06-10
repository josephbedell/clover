#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Genetic Fortune Decoded: Using LLMs to Discover Four-Leaf Clover Secrets"
subtitle.text = "A Journey into the Genetics of Luck with AI"

# Add the image only if it is in a supported format
image_path = 'four_leaf_clover.jpeg'
try:
    slide.shapes.add_picture(image_path, Inches(1), Inches(2.5), height=Inches(2))
except ValueError as e:
    print(f"Error adding picture: {e}")

# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = "Brief overview of the significance of four-leaf clovers.\nMention the role of genetics in determining the four-leaf trait.\nIntroduce the use of AI and LLMs in genetic research."

# Slide 3: Four-Leaf Clover and Genetics
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Four-Leaf Clover and Genetics"

# Slide 4: The Role of AI and LLMs
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "The Role of AI and LLMs"

# Slide 5: Identifying Candidate Genes
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Identifying Candidate Genes"
content.text = "- **Gene 1**: Description, location\n- **Gene 2**: Description, location\n- **Gene 3**: Description, location"

# Slide 6: Python Code Snippets
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Python Code Snippets"
content.text = "# Sample code to fetch gene sequences\nfrom Bio import Entrez, SeqIO\n\ndef fetch_sequence(gene_id):\n    Entrez.email = 'your.email@example.com'\n    handle = Entrez.efetch(db='nucleotide', id=gene_id, rettype='fasta', retmode='text')\n    return SeqIO.read(handle, 'fasta')"

# Slide 7: Data Analysis with AI
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Data Analysis with AI"

# Slide 8: Results and Findings
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Results and Findings"

# Slide 9: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = "Summarize findings with infographics.\nVisual representation of the genetic analysis process."

# Slide 10: Questions and Discussion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Questions and Discussion"
content.text = "Open floor for questions.\nEngage with the audience."

# Save the presentation
prs.save('four_leaf_clover_presentation.pptx')
