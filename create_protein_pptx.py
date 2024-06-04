#!/usr/bin/env python3

import logging
import sys
from Bio import SeqIO
from pptx import Presentation
from pptx.util import Inches

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def parse_genbank(file_path, start, end, output_file):
    """Parse the GenBank file and extract protein sequences within the specified range."""
    logging.info(f"Parsing GenBank file {file_path} and extracting proteins between {start}-{end}...")
    records = SeqIO.parse(file_path, "genbank")
    proteins = []

    for record in records:
        logging.info(f"Processing record: {record.id}")
        for feature in record.features:
            if feature.type == "CDS":
                location = feature.location
                logging.info(f"Found CDS feature from {location.start} to {location.end}")
                if (start <= location.start <= end) or (start <= location.end <= end):
                    if "translation" in feature.qualifiers:
                        protein_seq = feature.qualifiers["translation"][0]
                        protein_id = feature.qualifiers.get("protein_id", ["unknown_protein_id"])[0]
                        protein_desc = feature.qualifiers.get("product", ["unknown_product"])[0]
                        proteins.append((protein_id, protein_desc, protein_seq, location))
                        logging.info(f"Extracted protein: {protein_id} - {protein_desc}")

    with open(output_file, "w") as fasta_file:
        for protein_id, protein_desc, protein_seq, location in proteins:
            location_str = f"chr2:{location.start}..{location.end}"
            fasta_file.write(f">{protein_id} {protein_desc} {location_str}\n")
            fasta_file.write(f"{protein_seq}\n")

    if not proteins:
        logging.warning("No proteins found in the specified range.")

    return proteins

def create_presentation(proteins, pptx_file):
    """Create a PowerPoint presentation summarizing the extracted protein data."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Protein Summary"
    subtitle.text = "Extracted from GenBank file"

    # Protein slides
    for protein_id, protein_desc, protein_seq, location in proteins:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = f"{protein_id} - {protein_desc}"
        text_frame = body_shape.text_frame

        location_str = f"Location: chr2:{location.start}..{location.end}"
        seq_str = f"Sequence: {protein_seq[:60]}... (length: {len(protein_seq)} aa)"

        text_frame.text = location_str
        p = text_frame.add_paragraph()
        p.text = seq_str

    prs.save(pptx_file)
    logging.info(f"PowerPoint presentation saved to {pptx_file}")

def main():
    if len(sys.argv) < 3:
        print("Usage: ./create_protein_pptx.py <genbank_file> <output_pptx_file> [<start> <end>]")
        sys.exit(1)

    genbank_file_path = sys.argv[1]
    output_pptx_file = sys.argv[2]
    start = int(sys.argv[3]) if len(sys.argv) > 3 else 5000000
    end = int(sys.argv[4]) if len(sys.argv) > 4 else 10000000
    output_fasta_file = "chr2_region_proteins.faa"

    proteins = parse_genbank(genbank_file_path, start, end, output_fasta_file)
    if proteins:
        create_presentation(proteins, output_pptx_file)
    logging.info("Script completed.")

if __name__ == "__main__":
    logging.info("Starting protein extraction from GenBank file...")
    main()
    logging.info("Script completed.")
